#!/usr/bin/env python3
"""Build an architecture-defense slide deck (.pptx) from a deck.json spec.

Usage:
    pip install python-pptx
    python build_deck.py deck.json ArchDefense.pptx

`deck.json` is either a list of slide objects, or an object
`{ "brand": "PLUGFORGE", "slides": [ ... ] }`. Each slide object has a
`type` field; see SKILL.md for the full schema and deck.example.json for a
worked example covering every slide type.

This is a TEMPLATE: it owns the theme (colors, fonts, layout archetypes) so
the skill only has to produce structured content, never raw pptx geometry.
"""
import json
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- Theme -----------------------------------------------------------------
INK     = RGBColor(0x0B, 0x14, 0x22)   # deep navy: dark-slide bg / dark text
PANEL   = RGBColor(0x16, 0x22, 0x33)   # card/panel on dark slides
PANEL_L = RGBColor(0x2A, 0x38, 0x4A)   # hairline on dark slides
LIGHT   = RGBColor(0xF5, 0xF7, 0xFA)   # light-slide bg
CARD    = RGBColor(0xFF, 0xFF, 0xFF)   # white card on light slides
CARD_LN = RGBColor(0xE2, 0xE6, 0xEC)   # card border on light slides
TEAL    = RGBColor(0x33, 0xC9, 0xA8)   # primary accent
TEAL_BG = RGBColor(0xEC, 0xFB, 0xF7)   # faint teal fill for "why" boxes
ORANGE  = RGBColor(0xE8, 0xA3, 0x3D)   # internal / warning accent
RED     = RGBColor(0xE5, 0x56, 0x4E)   # before / danger accent
MUTED   = RGBColor(0x6B, 0x74, 0x84)   # muted body on light
MUTED_D = RGBColor(0x9A, 0xA6, 0xB2)   # muted body on dark
FAINT   = RGBColor(0x4A, 0x55, 0x63)   # very low-contrast footer text
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

MONO = "Consolas"      # display / headers / code (monospace, ships on Windows)
SANS = "Segoe UI"      # body copy

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)
DARK_TYPES = ("title", "pipeline", "summary")


# ---- Run / paragraph shorthands -------------------------------------------
def R(text, **spec):
    """A styled run: R('Hi', size=24, bold=True, color=WHITE)."""
    return (text, spec)


def P(*runs, **opts):
    """A paragraph made of runs, plus alignment / spacing options."""
    opts["runs"] = list(runs)
    return opts


# ---- Primitive helpers -----------------------------------------------------
def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _box(slide, left, top, width, height, fill=None, line=None,
         line_w=Pt(1), radius=True):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shape, int(left), int(top), int(width), int(height))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w
    sp.shadow.inherit = False
    return sp


def _text(slide, left, top, width, height, paragraphs,
          anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(para.get("space_after", 4))
        p.space_before = Pt(para.get("space_before", 0))
        if "line_spacing" in para:
            p.line_spacing = para["line_spacing"]
        for text, spec in para["runs"]:
            run = p.add_run()
            run.text = text
            f = run.font
            f.name = spec.get("font", SANS)
            f.size = Pt(spec.get("size", 14))
            f.bold = spec.get("bold", False)
            f.italic = spec.get("italic", False)
            f.color.rgb = spec.get("color", INK)
    return tb


def _header(slide, d, dark):
    """Kicker + title + subtitle. Returns the y where content can start."""
    title_color = WHITE if dark else INK
    sub_color = MUTED_D if dark else MUTED
    y = Inches(0.55)
    if d.get("kicker"):
        _text(slide, Inches(0.6), y, Inches(12.1), Inches(0.3),
              [P(R(d["kicker"].upper(), font=MONO, size=11, bold=True, color=TEAL))])
        y = Inches(0.92)
    _text(slide, Inches(0.6), y, Inches(12.1), Inches(0.8),
          [P(R(d.get("title", ""), font=MONO, size=30, bold=True, color=title_color))])
    y = y + Inches(0.82)
    if d.get("subtitle"):
        _text(slide, Inches(0.6), y, Inches(12.1), Inches(0.6),
              [P(R(d["subtitle"], size=14, color=sub_color), line_spacing=1.1)])
        y = y + Inches(0.62)
    return y + Inches(0.15)


def _footer(slide, idx, brand, dark):
    color = MUTED_D if dark else MUTED
    _text(slide, Inches(0.6), Inches(7.06), Inches(8), Inches(0.3),
          [P(R(brand.upper(), font=MONO, size=9, bold=True, color=color))])
    _text(slide, Inches(11.4), Inches(7.06), Inches(1.3), Inches(0.3),
          [P(R(f"{idx:02d}", font=MONO, size=9, color=color), align=PP_ALIGN.RIGHT)])


# ---- Slide builders --------------------------------------------------------
def slide_title(slide, d):
    _bg(slide, INK)
    _box(slide, 0, 0, Inches(0.18), EMU_H, fill=TEAL, radius=False)
    _text(slide, Inches(0.9), Inches(1.6), Inches(11), Inches(0.4),
          [P(R(d.get("kicker", "ARCHITECTURE DEFENSE").upper(),
               font=MONO, size=13, bold=True, color=TEAL))])
    _text(slide, Inches(0.85), Inches(2.05), Inches(11.6), Inches(1.5),
          [P(R(d.get("title", ""), font=MONO, size=58, bold=True, color=WHITE))])
    if d.get("subtitle"):
        _text(slide, Inches(0.9), Inches(3.75), Inches(11), Inches(0.6),
              [P(R(d["subtitle"], size=19, color=MUTED_D))])
    if d.get("terminal"):
        _box(slide, Inches(0.9), Inches(4.55), Inches(9.6), Inches(0.72),
             fill=PANEL, line=TEAL)
        _text(slide, Inches(1.15), Inches(4.55), Inches(9.2), Inches(0.72),
              [P(R(d["terminal"], font=MONO, size=15, color=TEAL))],
              anchor=MSO_ANCHOR.MIDDLE)
    if d.get("thesis"):
        _text(slide, Inches(0.9), Inches(5.55), Inches(11.2), Inches(0.9),
              [P(R(d["thesis"], italic=True, size=13, color=MUTED_D), line_spacing=1.2)])
    if d.get("footer"):
        _text(slide, Inches(0.9), Inches(6.65), Inches(11), Inches(0.4),
              [P(R(d["footer"], font=MONO, size=11, color=FAINT))])


def slide_cards(slide, d):
    _bg(slide, LIGHT)
    top = _header(slide, d, dark=False)
    cards = d.get("cards", [])
    n = max(1, len(cards))
    gap, left, total_w = Inches(0.3), Inches(0.6), Inches(12.1)
    cw = int((total_w - gap * (n - 1)) / n)
    ch = Inches(2.7)
    for i, c in enumerate(cards):
        x = int(left + i * (cw + gap))
        _box(slide, x, top, cw, ch, fill=CARD, line=CARD_LN)
        _box(slide, x, top, cw, Inches(0.08), fill=TEAL, radius=False)
        _text(slide, x + Inches(0.25), top + Inches(0.32), cw - Inches(0.5), Inches(0.7),
              [P(R(c.get("title", ""), font=MONO, size=15, bold=True, color=INK),
                 line_spacing=1.05)])
        _text(slide, x + Inches(0.25), top + Inches(1.05), cw - Inches(0.5), ch - Inches(1.2),
              [P(R(c.get("body", ""), size=12, color=MUTED), line_spacing=1.18)])
    if d.get("callout"):
        co = d["callout"]
        cy = Inches(5.95)
        _box(slide, Inches(0.6), cy, Inches(12.1), Inches(0.95), fill=INK)
        _box(slide, Inches(0.6), cy, Inches(0.08), Inches(0.95), fill=TEAL, radius=False)
        _text(slide, Inches(0.95), cy + Inches(0.13), Inches(11.5), Inches(0.3),
              [P(R(co.get("title", "").upper(), font=MONO, size=11, bold=True, color=TEAL))])
        _text(slide, Inches(0.95), cy + Inches(0.43), Inches(11.5), Inches(0.45),
              [P(R(co.get("body", ""), size=13, color=WHITE), line_spacing=1.1)])


def slide_pipeline(slide, d):
    _bg(slide, INK)
    top = _header(slide, d, dark=True)
    steps = d.get("steps", [])
    n = max(1, len(steps))
    hl = d.get("highlight", -1)
    gap, left, total_w = Inches(0.45), Inches(0.6), Inches(12.1)
    bw = int((total_w - gap * (n - 1)) / n)
    bh = Inches(1.45)
    for i, s in enumerate(steps):
        x = int(left + i * (bw + gap))
        fill = TEAL if i == hl else PANEL
        txt = INK if i == hl else WHITE
        _box(slide, x, top, bw, bh, fill=fill,
             line=(None if i == hl else PANEL_L))
        _text(slide, x, top, bw, bh,
              [P(R(s, font=MONO, size=13, bold=True, color=txt),
                 align=PP_ALIGN.CENTER, line_spacing=1.05)],
              anchor=MSO_ANCHOR.MIDDLE)
        if i < n - 1:
            _text(slide, int(x + bw), top, int(gap), bh,
                  [P(R("›", font=MONO, size=20, bold=True, color=TEAL),
                     align=PP_ALIGN.CENTER)],
                  anchor=MSO_ANCHOR.MIDDLE)
    panels = d.get("panels", [])
    if panels:
        py = int(top + bh + Inches(0.5))
        pn = len(panels)
        pgap = Inches(0.4)
        pw = int((total_w - pgap * (pn - 1)) / pn)
        ph = Inches(2.5)
        for j, pnl in enumerate(panels):
            x = int(left + j * (pw + pgap))
            _box(slide, x, py, pw, ph, fill=PANEL)
            hc = ORANGE if pnl.get("accent") == "orange" else TEAL
            _text(slide, x + Inches(0.3), py + Inches(0.2), pw - Inches(0.6), Inches(0.4),
                  [P(R(pnl.get("title", ""), font=MONO, size=14, bold=True, color=hc))])
            bullets = [P(R("•  " + b, size=12, color=MUTED_D),
                         line_spacing=1.1, space_after=6)
                       for b in pnl.get("bullets", [])]
            _text(slide, x + Inches(0.3), py + Inches(0.7), pw - Inches(0.6), ph - Inches(0.9),
                  bullets)


def slide_two_column(slide, d):
    _bg(slide, LIGHT)
    top = _header(slide, d, dark=False)
    cols = d.get("columns", [])
    n = max(1, len(cols))
    gap, left, total_w = Inches(0.5), Inches(0.6), Inches(12.1)
    cw = int((total_w - gap * (n - 1)) / n)
    total_h = Inches(4.5)
    for i, col in enumerate(cols):
        x = int(left + i * (cw + gap))
        accent = RED if col.get("accent") == "red" else TEAL
        _box(slide, x, top, cw, Inches(0.7), fill=INK)
        _text(slide, x + Inches(0.3), top, cw - Inches(0.6), Inches(0.7),
              [P(R(col.get("heading", ""), font=MONO, size=16, bold=True, color=WHITE),
                 R(("   " + col["subhead"]) if col.get("subhead") else "",
                   size=11, color=accent))],
              anchor=MSO_ANCHOR.MIDDLE)
        body_top = int(top + Inches(0.7))
        body_h = int(total_h - Inches(0.7))
        _box(slide, x, body_top, cw, body_h, fill=CARD, line=CARD_LN)
        yy = body_top + Inches(0.28)
        for k, step in enumerate(col.get("steps", [])):
            _box(slide, x + Inches(0.28), yy, Inches(0.34), Inches(0.34), fill=accent)
            _text(slide, x + Inches(0.28), yy, Inches(0.34), Inches(0.34),
                  [P(R(str(k + 1), font=MONO, size=11, bold=True, color=WHITE),
                     align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
            _text(slide, x + Inches(0.75), yy - Inches(0.03), cw - Inches(1.05), Inches(0.5),
                  [P(R(step, size=12, color=INK), line_spacing=1.05)])
            yy = yy + Inches(0.52)
        if col.get("why"):
            wy = int(body_top + body_h - Inches(1.05))
            _box(slide, x + Inches(0.22), wy, cw - Inches(0.44), Inches(0.85),
                 fill=TEAL_BG, line=accent)
            _text(slide, x + Inches(0.38), wy + Inches(0.1), cw - Inches(0.76), Inches(0.65),
                  [P(R("Why: ", bold=True, size=11, color=accent),
                     R(col["why"], size=11, color=INK), line_spacing=1.12)])


def slide_metrics(slide, d):
    _bg(slide, LIGHT)
    top = _header(slide, d, dark=False)
    metrics = d.get("metrics", [])
    if d.get("terminal"):
        _box(slide, Inches(0.6), top, Inches(7.0), Inches(3.4), fill=INK)
        lines = d["terminal"].split("\n")
        paras = []
        for ln in lines:
            hot = ("✓" in ln) or ln.strip().startswith("→")
            paras.append(P(R(ln, font=MONO, size=14,
                             color=(TEAL if hot else WHITE)), space_after=6))
        _text(slide, Inches(0.9), top + Inches(0.3), Inches(6.4), Inches(2.9), paras)
        mx, mw = Inches(7.9), Inches(4.8)
    else:
        mx, mw = Inches(0.6), Inches(12.1)
    my = top
    for m in metrics:
        _box(slide, mx, my, mw, Inches(0.72), fill=PANEL)
        _text(slide, mx + Inches(0.3), my, Inches(1.9), Inches(0.72),
              [P(R(m.get("value", ""), font=MONO, size=22, bold=True, color=TEAL))],
              anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, mx + Inches(2.3), my, mw - Inches(2.6), Inches(0.72),
              [P(R(m.get("label", ""), size=12, color=WHITE))],
              anchor=MSO_ANCHOR.MIDDLE)
        my = my + Inches(0.86)
    if d.get("callout"):
        cy = Inches(5.95)
        _box(slide, Inches(0.6), cy, Inches(12.1), Inches(0.95), fill=INK, line=TEAL)
        _text(slide, Inches(0.95), cy, Inches(11.5), Inches(0.95),
              [P(R(d["callout"], size=13, color=WHITE), line_spacing=1.15)],
              anchor=MSO_ANCHOR.MIDDLE)


def slide_timeline(slide, d):
    _bg(slide, LIGHT)
    top = _header(slide, d, dark=False)
    items = d.get("items", [])
    yy = top
    rh = Inches(0.62)
    for it in items:
        _box(slide, Inches(0.6), yy, Inches(0.5), Inches(0.5), fill=TEAL)
        _text(slide, Inches(0.6), yy, Inches(0.5), Inches(0.5),
              [P(R(it.get("badge", ""), font=MONO, size=12, bold=True, color=WHITE),
                 align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, Inches(1.3), yy, Inches(3.7), Inches(0.5),
              [P(R(it.get("title", ""), font=MONO, size=15, bold=True, color=INK))],
              anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, Inches(5.2), yy, Inches(7.5), Inches(0.5),
              [P(R(it.get("body", ""), size=12, color=MUTED), line_spacing=1.05)],
              anchor=MSO_ANCHOR.MIDDLE)
        yy = yy + rh


def slide_summary(slide, d):
    _bg(slide, INK)
    top = _header(slide, d, dark=True)
    rows = d.get("rows", [])
    yy = top
    rh = Inches(0.76)
    for r in rows:
        _box(slide, Inches(0.6), yy, Inches(12.1), Inches(0.64), fill=PANEL)
        _box(slide, Inches(0.6), yy, Inches(0.08), Inches(0.64), fill=TEAL, radius=False)
        _text(slide, Inches(1.0), yy, Inches(4.3), Inches(0.64),
              [P(R(r.get("title", ""), font=MONO, size=14, bold=True, color=WHITE))],
              anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, Inches(5.4), yy, Inches(7.1), Inches(0.64),
              [P(R(r.get("body", ""), size=12, color=MUTED_D), line_spacing=1.05)],
              anchor=MSO_ANCHOR.MIDDLE)
        yy = yy + rh
    if d.get("footer"):
        _text(slide, Inches(0.6), Inches(6.65), Inches(12.1), Inches(0.5),
              [P(R(d["footer"], font=MONO, size=14, bold=True, italic=True, color=TEAL),
                 align=PP_ALIGN.CENTER)])


def slide_image(slide, d):
    dark = bool(d.get("dark"))
    _bg(slide, INK if dark else LIGHT)
    top = _header(slide, d, dark=dark)
    img = d.get("image")
    if img:
        try:
            pic = slide.shapes.add_picture(img, Inches(0.9), top)
            max_w = Inches(11.5)
            avail_h = int(Inches(6.95) - top)
            scale = min(max_w / pic.width, avail_h / pic.height, 1.0)
            pic.width = int(pic.width * scale)
            pic.height = int(pic.height * scale)
            pic.left = int((EMU_W - pic.width) / 2)
            pic.top = top
        except Exception as exc:  # noqa: BLE001 - placeholder beats a crash
            _text(slide, Inches(0.9), top + Inches(1.5), Inches(11.5), Inches(1),
                  [P(R(f"[ image not found: {img} — {exc} ]", font=MONO,
                       size=14, color=RED), align=PP_ALIGN.CENTER)])
    if d.get("caption"):
        _text(slide, Inches(0.6), Inches(7.0), Inches(12.1), Inches(0.35),
              [P(R(d["caption"], italic=True, size=11,
                   color=(MUTED_D if dark else MUTED)), align=PP_ALIGN.CENTER)])


BUILDERS = {
    "title": slide_title,
    "cards": slide_cards,
    "pipeline": slide_pipeline,
    "two_column": slide_two_column,
    "metrics": slide_metrics,
    "timeline": slide_timeline,
    "summary": slide_summary,
    "image": slide_image,
}


def build(slides, outfile, brand="ARCH DEFENSE"):
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    blank = prs.slide_layouts[6]
    for idx, d in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank)
        t = d.get("type")
        builder = BUILDERS.get(t)
        if builder is None:
            raise SystemExit(f"Unknown slide type {t!r} on slide {idx}. "
                             f"Valid types: {', '.join(sorted(BUILDERS))}")
        builder(slide, d)
        if t != "title":
            _footer(slide, idx, d.get("brand", brand), dark=(t in DARK_TYPES))
    prs.save(outfile)
    print(f"Wrote {outfile} ({len(slides)} slides)")


def main():
    if len(sys.argv) < 3:
        print("usage: python build_deck.py deck.json out.pptx")
        sys.exit(1)
    with open(sys.argv[1], encoding="utf-8") as fh:
        data = json.load(fh)
    if isinstance(data, dict):
        brand = data.get("brand", "ARCH DEFENSE")
        slides = data["slides"]
    else:
        brand = "ARCH DEFENSE"
        slides = data
    build(slides, sys.argv[2], brand=brand)


if __name__ == "__main__":
    main()
